#!/usr/bin/env python3
"""
Enhanced File Sorter for Cybersecurity Portfolio
Automatically sorts documents into the appropriate portfolio categories
"""

import os
import shutil
from pathlib import Path
from docx import Document
from pptx import Presentation

# Define keywords for each portfolio category
CATEGORIES = {
    "incident_response": [
        "incident", "response", "handler", "journal", "phishing", "alert", "ticket",
        "usb", "security", "exercise", "breach", "threat", "detection"
    ],
    "vulnerability_assessment": [
        "vulnerability", "assessment", "scan", "penetration", "testing", "sql", "injection",
        "risk", "security", "test", "analysis", "report", "finding"
    ],
    "access_control": [
        "access", "control", "rbac", "identity", "authentication", "authorization",
        "permission", "linux", "file", "data", "leak", "prevention"
    ],
    "network_security": [
        "network", "monitoring", "traffic", "packet", "wireshark", "tcpdump",
        "device", "management", "protocol", "firewall", "ids", "ips"
    ],
    "tools_automation": [
        "automation", "script", "tool", "python", "programming", "development",
        "vps", "update", "backup", "monitoring", "security", "toolkit"
    ],
    "educational": [
        "lesson", "module", "summary", "explained", "training", "learning",
        "comparison", "analysis", "methodology", "framework", "study"
    ],
    "certifications": [
        "certification", "certificate", "credential", "badge", "completion",
        "course", "training", "accreditation"
    ],
    "resume": [
        "resume", "cv", "curriculum", "vitae", "professional", "experience"
    ]
}

def get_category(filename, content):
    """Determine the best category for a file based on filename and content"""
    filename_lower = filename.lower()
    content_lower = content.lower()
    
    # Score each category based on keyword matches
    scores = {}
    for category, keywords in CATEGORIES.items():
        score = 0
        for keyword in keywords:
            if keyword in filename_lower:
                score += 2  # Higher weight for filename matches
            if keyword in content_lower:
                score += 1  # Lower weight for content matches
        scores[category] = score
    
    # Return the category with the highest score
    if scores:
        best_category = max(scores, key=scores.get)
        if scores[best_category] > 0:
            return best_category
    
    return None

def read_docx(file_path):
    """Extract text content from DOCX files"""
    try:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Warning: Could not read DOCX content from {file_path}: {e}")
        return ""

def read_pptx(file_path):
    """Extract text content from PPTX files"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Warning: Could not read PPTX content from {file_path}: {e}")
        return ""

def sort_files():
    """Main function to sort files into portfolio categories"""
    current_dir = Path.cwd()
    portfolio_dir = current_dir / "Portfolio" / "projects"
    
    # Check if we're in the right directory
    if not portfolio_dir.exists():
        print("❌ Error: Portfolio directory not found!")
        print("Please run this script from the Cybersecurity folder.")
        return
    
    # Get all document files in current directory
    document_extensions = ('.docx', '.pptx', '.pdf', '.txt', '.md')
    files = [f for f in current_dir.iterdir() 
             if f.is_file() and f.suffix.lower() in document_extensions]
    
    if not files:
        print("✅ No document files found to sort.")
        return
    
    print("🛡️  Cybersecurity Portfolio File Sorter")
    print("=" * 50)
    print(f"Found {len(files)} document(s) to sort...")
    print()
    
    moved_files = []
    skipped_files = []
    
    for file_path in files:
        print(f"📄 Processing: {file_path.name}")
        
        # Extract content based on file type
        content = ""
        if file_path.suffix.lower() == '.docx':
            content = read_docx(file_path)
        elif file_path.suffix.lower() == '.pptx':
            content = read_pptx(file_path)
        
        # Determine category
        category = get_category(file_path.name, content)
        
        if category:
            # Create destination path
            dest_dir = portfolio_dir / category
            dest_dir.mkdir(exist_ok=True)
            dest_path = dest_dir / file_path.name
            
            # Move file
            try:
                shutil.move(str(file_path), str(dest_path))
                moved_files.append((file_path.name, category))
                print(f"✅ Moved to: {category}/")
            except Exception as e:
                print(f"❌ Error moving file: {e}")
                skipped_files.append(file_path.name)
        else:
            print(f"❓ No category found - file left in place")
            skipped_files.append(file_path.name)
        
        print()
    
    # Summary
    print("=" * 50)
    print("📊 Sorting Summary:")
    print(f"✅ Successfully moved: {len(moved_files)} files")
    print(f"❓ Skipped: {len(skipped_files)} files")
    print()
    
    if moved_files:
        print("📁 Files moved:")
        for filename, category in moved_files:
            print(f"  • {filename} → {category}/")
        print()
    
    if skipped_files:
        print("📄 Files that need manual sorting:")
        for filename in skipped_files:
            print(f"  • {filename}")
        print()
        print("💡 Tip: Use the portfolio automation tool to add these manually:")
        print("  python3 run_portfolio.py")
    
    print("🎉 File sorting complete!")
    print("Next step: Use 'python3 run_portfolio.py' to add files to your portfolio with documentation.")

def main():
    """Entry point"""
    try:
        sort_files()
    except KeyboardInterrupt:
        print("\n\n👋 File sorting cancelled.")
    except Exception as e:
        print(f"\n❌ Error during file sorting: {e}")

if __name__ == "__main__":
    main()